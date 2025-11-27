import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import io
from PIL import Image

class PDFToPPTConverter:
    def __init__(self, pdf_file):
        """
        Initialize with a PDF file (stream or path).
        """
        self.pdf_file = pdf_file
        self.doc = fitz.open(stream=pdf_file.read(), filetype="pdf")

    def convert_to_images(self, dpi=200):
        """
        Default Mode: Convert each page to an image and place on a slide.
        Returns a BytesIO object containing the PPTX file.
        """
        prs = Presentation()
        
        # Set slide dimensions to match the first page of PDF (assuming uniform size)
        if len(self.doc) > 0:
            page = self.doc[0]
            # PyMuPDF uses points (1/72 inch), pptx uses EMUs (914400 per inch)
            width_inch = page.rect.width / 72
            height_inch = page.rect.height / 72
            
            prs.slide_width = int(width_inch * 914400)
            prs.slide_height = int(height_inch * 914400)

        for page_num in range(len(self.doc)):
            page = self.doc[page_num]
            
            # Render page to image (high resolution)
            # Default PDF is 72 DPI. Zoom = dpi / 72
            zoom = dpi / 72
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            
            # Convert to PIL Image
            img_data = pix.tobytes("png")
            image_stream = BytesIO(img_data)
            
            # Add slide
            blank_slide_layout = prs.slide_layouts[6] # 6 is blank
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to slide
            slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    def convert_separated(self):
        """
        Separation Mode: Extract images and text.
        Returns a BytesIO object containing the PPTX file.
        """
        prs = Presentation()
        
        # Set slide dimensions (standard 16:9 or based on PDF?)
        # Let's stick to PDF dimensions for consistency
        if len(self.doc) > 0:
            page = self.doc[0]
            width_inch = page.rect.width / 72
            height_inch = page.rect.height / 72
            prs.slide_width = int(width_inch * 914400)
            prs.slide_height = int(height_inch * 914400)

        for page_num in range(len(self.doc)):
            page = self.doc[page_num]
            
            # Create a slide for this page
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 1. Extract and place images
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = self.doc.extract_image(xref)
                    optimized_bytes = self._optimize_image_bytes(base_image["image"])
                    rects = page.get_image_rects(xref)

                    if rects:
                        for rect in rects:
                            left = Inches(rect.x0 / 72)
                            top = Inches(rect.y0 / 72)
                            width = Inches(rect.width / 72)
                            height = Inches(rect.height / 72)
                            if width <= 0 or height <= 0:
                                continue
                            slide.shapes.add_picture(
                                BytesIO(optimized_bytes),
                                left,
                                top,
                                width=width,
                                height=height
                            )
                    else:
                        fallback_left = Inches(0.5 + (img_index % 2) * 3)
                        fallback_top = Inches(0.5 + (img_index // 2) * 3.5)
                        slide.shapes.add_picture(
                            BytesIO(optimized_bytes),
                            fallback_left,
                            fallback_top,
                            height=Inches(3)
                        )
                except Exception as image_error:
                    print(f"Could not place image {img_index}: {image_error}")
                    continue

            # 2. Extract and place text
            # get_text("dict") gives us blocks with bbox
            text_blocks = page.get_text("dict")["blocks"]
            for block_index, block in enumerate(text_blocks):
                if block.get("type") != 0:
                    continue
                bbox = block.get("bbox")
                if not bbox:
                    continue

                left = Inches(bbox[0] / 72)
                top = Inches(bbox[1] / 72)
                width = max(Inches((bbox[2] - bbox[0]) / 72), Inches(0.25))
                height = max(Inches((bbox[3] - bbox[1]) / 72), Inches(0.25))

                try:
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                except Exception as text_box_error:
                    print(f"Could not add textbox for block {block_index}: {text_box_error}")
                    continue

                text_frame = text_box.text_frame
                text_frame.clear()
                text_frame.word_wrap = True

                for line in block.get("lines", []):
                    if not line.get("spans"):
                        continue
                    paragraph = text_frame.add_paragraph()

                    for span in line["spans"]:
                        run = paragraph.add_run()
                        run.text = span.get("text", "")

                        try:
                            run.font.size = Pt(span.get("size", 12))

                            color_int = span.get("color", 0)
                            r = (color_int >> 16) & 0xFF
                            g = (color_int >> 8) & 0xFF
                            b = color_int & 0xFF
                            run.font.color.rgb = RGBColor(r, g, b)

                            flags = span.get("flags", 0)
                            if flags & 2**4:
                                run.font.bold = True
                            if flags & 2**1:
                                run.font.italic = True
                        except Exception:
                            continue

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    def _optimize_image_bytes(self, image_bytes, max_dimension=2000):
        """Downscale and compress images to reduce memory usage inside PPT."""
        with Image.open(BytesIO(image_bytes)) as image:
            if image.mode not in ("RGB", "L"):
                image = image.convert("RGB")
            image.thumbnail((max_dimension, max_dimension), Image.LANCZOS)
            optimized_stream = BytesIO()
            image.save(optimized_stream, format="PNG", optimize=True)
            optimized_stream.seek(0)
            return optimized_stream.getvalue()

    def extract_text_content(self):
        """
        Extract all text from the PDF.
        Returns a BytesIO object containing the text file.
        """
        text_output = io.StringIO()
        
        for page_num in range(len(self.doc)):
            page = self.doc[page_num]
            text = page.get_text()
            text_output.write(f"--- Page {page_num + 1} ---\n\n")
            text_output.write(text)
            text_output.write("\n\n")
            
        output = BytesIO(text_output.getvalue().encode('utf-8'))
        output.seek(0)
        return output
